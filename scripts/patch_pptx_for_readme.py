"""Patch presentation slides for README export (spacing, labels, bullets)."""
from __future__ import annotations

import shutil
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "presentation_source.pptx"

DEFAULT_FONT = "Pretendard Medium"
DEFAULT_BOLD = "Pretendard Bold"


def _clear_bullets(paragraph) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag in (
            qn("a:buChar"),
            qn("a:buAutoNum"),
            qn("a:buBlip"),
            qn("a:buFont"),
            qn("a:buClr"),
            qn("a:buSzPct"),
            qn("a:buSzPts"),
        ):
            p_pr.remove(child)


def _read_font(paragraph) -> tuple[str, object | None]:
    if paragraph.runs:
        run = paragraph.runs[0]
        return run.font.name or DEFAULT_FONT, run.font.size
    return DEFAULT_FONT, Pt(12)


def _apply_font(tf, name: str, size) -> None:
    for para in tf.paragraphs:
        _clear_bullets(para)
        para.level = 0
        for run in para.runs:
            run.font.name = name
            if size is not None:
                run.font.size = size


def _write_single_line(shape, text: str, *, font_name: str | None = None, font_size=None) -> None:
    tf = shape.text_frame
    name, size = _read_font(tf.paragraphs[0])
    if font_name:
        name = font_name
    if font_size is not None:
        size = font_size

    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.text = text
    _apply_font(tf, name, size)


def _write_multiline(shape, lines: list[str]) -> None:
    tf = shape.text_frame
    name, size = _read_font(tf.paragraphs[0])
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.text = "\n".join(lines)
    _apply_font(tf, name, size)


def _no_wrap(shape, min_width_in: float) -> None:
    shape.text_frame.word_wrap = False
    shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    if shape.width < int(Inches(min_width_in)):
        shape.width = int(Inches(min_width_in))


def _plain_bullet(shape, text: str) -> None:
    tf = shape.text_frame
    tf.margin_left = Pt(18)
    tf.margin_right = Pt(10)
    clean = text.lstrip("• ").strip()
    _write_single_line(shape, f"• {clean}")


def _patch_hero(slide) -> None:
    to_remove = []
    sympo_shape = None
    title_shape = None

    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and sh.top > int(Inches(7.5)):
            to_remove.append(sh._element)
            continue
        if not sh.has_text_frame:
            continue
        text = sh.text_frame.text.strip()
        if text == "A3" or "박선민" in text:
            to_remove.append(sh._element)
        elif text == "SYMPO":
            sympo_shape = sh
        elif "Multi-Agent Orchestration" in text:
            title_shape = sh

    for el in to_remove:
        el.getparent().remove(el)

    if sympo_shape is not None:
        sympo_shape.width = int(Inches(11.2))
        sympo_shape.text_frame.word_wrap = False

    if title_shape is not None:
        title_shape.left = int(Inches(12.8))
        title_shape.width = int(Inches(13.8))
        title_shape.text_frame.word_wrap = False


def _patch_pipeline(slide) -> None:
    # Micro-labels overlaid directly on the illustration icons. Each stage's
    # bottom caption already names it, so these overlays are redundant clutter
    # that collides with the graphics ("글자 깨짐"). Drop them for a clean
    # icon + stage-title + caption layout.
    drop_labels = {
        "회의 음성",
        "프로젝트\n명세서",
        "WBS 생성\n에이전트",
        "태스크 관리\n에이전트",
        "팀원\n에이전트",
    }
    caption_lines = {
        "태스크별 성향\u00a0기반\n팀원\u00a0에이전트 호출": ["태스크별 성향 기반", "팀원 에이전트 호출"],
        "계층형 프로젝트\u00a0\nWBS 초안 생성": ["계층형 프로젝트", "WBS 초안 생성"],
        "회의 음성 및\u00a0\n프로젝트 명세서 입력": ["회의 음성 및", "프로젝트 명세서 입력"],
        "팀원 에이전트 간\n토의 후 피드백 요청": ["팀원 에이전트 간", "토의 후 피드백 요청"],
        "피드백 반영 후 업무\u00a0\n분배된 최종 WBS 도출": ["피드백 반영 후", "최종 WBS 도출"],
    }

    to_remove = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        text = sh.text_frame.text
        if text in drop_labels:
            to_remove.append(sh._element)
        elif text in caption_lines:
            _write_multiline(sh, caption_lines[text])
            sh.width = int(Inches(4.65))
            sh.text_frame.margin_left = Pt(6)
            sh.text_frame.margin_right = Pt(6)

    for el in to_remove:
        el.getparent().remove(el)

    # Debate speech labels are boxed too narrow and wrap mid-word
    # ("수정 요\n청!"). Keep them on a single line.
    debate_labels = {"수정 요청!", "수정 명령!"}
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text in debate_labels:
            sh.text_frame.word_wrap = False
            sh.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            sh.width = int(Inches(1.7))

    for sh in slide.shapes:
        if sh.has_text_frame and "팀원 메타 데이터" in sh.text_frame.text:
            _no_wrap(sh, 20.0)


def _patch_outputs(slide) -> None:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        text = sh.text_frame.text.strip()
        if not text:
            continue

        if text == "프로젝트 출력물":
            _write_single_line(sh, text, font_name=DEFAULT_BOLD)
            _no_wrap(sh, 8.5)
            continue

        if text in ("WBS 최종본(작업 분할 구조도)", "프로젝트 일정표 (Gantt Chart)"):
            _no_wrap(sh, 10.0)
            continue

        if text in ("WBS 최종본 산출 구조", "일정 시각화 특징"):
            _write_single_line(sh, text, font_name=DEFAULT_BOLD)
            _no_wrap(sh, 6.2)
            sh.height = int(Inches(0.85))
            sh.text_frame.margin_bottom = Pt(10)
            continue

        if text.startswith(("3단계", "R&R", "프로젝트 메트릭", "직관적", "태스크 연결성")):
            _plain_bullet(sh, text)
            sh.width = int(Inches(11.8))
            sh.top += int(Inches(0.12))


def patch_presentation(src: Path = SOURCE) -> Path:
    if not src.exists():
        raise FileNotFoundError(src)

    tmp = Path(tempfile.mkdtemp(prefix="sympo_pptx_"))
    out = tmp / "readme_export.pptx"
    shutil.copy2(src, out)

    prs = Presentation(str(out))
    _patch_hero(prs.slides[0])
    _patch_pipeline(prs.slides[9])
    _patch_outputs(prs.slides[35])
    prs.save(str(out))
    return out


if __name__ == "__main__":
    path = patch_presentation()
    print("patched", path)
