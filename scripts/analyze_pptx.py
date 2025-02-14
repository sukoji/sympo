"""Analyze presentation slides and export hero images."""
from __future__ import annotations

import os
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "presentation_source.pptx"
OUT = ROOT / "docs" / "assets"
EXTRACT = OUT / "pptx_extract"


def main():
    prs = Presentation(str(PPTX))
    lines = []
    print("slides:", len(prs.slides))
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        pic_count = 0
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() and not title:
                title = shape.text_frame.text.strip().split("\n")[0][:100]
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pic_count += 1
        lines.append(f"Slide {i:2d} | pics={pic_count:2d} | {title}")
    out = ROOT / "docs" / "assets" / "slide_index.txt"
    out.write_text("\n".join(lines), encoding="utf-8")
    print("wrote", out)


if __name__ == "__main__":
    main()
