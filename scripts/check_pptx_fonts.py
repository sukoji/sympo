"""Warn if fonts used in presentation_source.pptx are missing on this machine."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "presentation_source.pptx"


def _installed_fonts() -> set[str]:
    import platform

    fonts: set[str] = set()
    if platform.system() == "Windows":
        import winreg

        for hive in (winreg.HKEY_LOCAL_MACHINE, winreg.HKEY_CURRENT_USER):
            try:
                key = winreg.OpenKey(hive, r"SOFTWARE\Microsoft\Windows NT\CurrentVersion\Fonts")
                i = 0
                while True:
                    try:
                        entry = winreg.EnumValue(key, i)
                        name = entry[0]
                        fonts.add(name.rsplit(" (", 1)[0])
                        i += 1
                    except OSError:
                        break
            except OSError:
                pass
        if fonts:
            return fonts

    try:
        from matplotlib import font_manager

        return {f.name for f in font_manager.fontManager.ttflist}
    except Exception:
        return fonts


def _pptx_fonts() -> set[str]:
    prs = Presentation(str(PPTX))
    fonts: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts.add(run.font.name)
    return fonts


def _font_available(name: str, installed: set[str]) -> bool:
    if name in installed:
        return True
    family = name.split()[0]
    return any(f == family or f.startswith(family + " ") for f in installed)


def main() -> int:
    if not PPTX.exists():
        print(f"missing: {PPTX}", file=sys.stderr)
        return 1

    installed = _installed_fonts()
    required = sorted(_pptx_fonts())
    missing = [f for f in required if not _font_available(f, installed)]
    critical = [f for f in missing if f.startswith("Pretendard")]
    optional = [f for f in missing if not f.startswith("Pretendard")]

    print("PPTX fonts:", ", ".join(required))
    if optional:
        print("\nOptional fonts missing (minor impact):")
        for f in optional:
            print(f"  - {f}")
    if critical:
        print("\nMISSING (install before export or text/layout will break):")
        for f in critical:
            print(f"  - {f}")
        print("\nInstall Pretendard: https://github.com/orioncactus/pretendard/releases")
        return 1

    print("\nPretendard OK. Run: python scripts/export_pptx_slides.py")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
